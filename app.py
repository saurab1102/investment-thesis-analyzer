import streamlit as st
from pptx import Presentation
from llm_client import query_together
from report_generator import generate_pdf_report
import json
import re
from datetime import datetime

REQUIRED_TYPES = [
    "Problem", "Solution", "Market", "Business Model",
    "Competition", "Team", "Financials", "Traction", "Funding Ask"
]

def extract_slide_text(file):
    prs = Presentation(file)
    slide_texts = []
    for i, slide in enumerate(prs.slides):
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    content.append(text)
        slide_texts.append({"slide_num": i + 1, "text": " ".join(content)})
    return slide_texts

def classify_single_slide(text):
    prompt = (
        "Classify the following slide into one of the categories:\n"
        "Problem, Solution, Market, Business Model, Competition, Team, Financials, Traction, Funding Ask, Unclassified.\n\n"
        f"Slide Content:\n\"{text}\"\n\n"
        "Return only the category as a single word."
    )
    return query_together(prompt)

def classify_all_slides(slides):
    for slide in slides:
        try:
            label = classify_single_slide(slide["text"])
            slide["category"] = label.strip()
        except Exception:
            slide["category"] = "Unclassified"
    return slides

def build_analysis_prompt(slides):
    classified = [s for s in slides if s.get("category") != "Unclassified"]
    header = (
        "You are an AI investment analyst. Based on the following pitch deck content, produce an investment thesis. "
        "Return your response as a JSON object with the following fields:\n\n"
        "- recommendation: one of [\"Strong Buy\", \"Hold\", \"Pass\"]\n"
        "- overall_score: integer 0–100\n"
        "- processing_date: use current UTC in format DD-MM-YYYY HH:MM:SS UTC\n"
        "- confidence_score: integer 0–100\n"
        "- strengths: list of 3–5 strings\n"
        "- weaknesses: list of 3–5 strings\n"
        "- recommendations: string (100–200 words)\n"
        "- categories: list of 9 objects, each with:\n"
        "  - name (category name),\n"
        "  - score (0–10),\n"
        "  - weight (int %),\n"
        "  - feedback (50–150 words)\n\n"
        "Use these fixed weights: Problem 10, Solution 15, Market 20, Business Model 15, Competition 10, Team 15, "
        "Traction 10, Financials 10, Clarity 5\n\n"
        "Classified Slides:\n"
    )
    grouped = {}
    for slide in classified:
        cat = slide["category"]
        grouped.setdefault(cat, []).append(slide["text"])
    body = ""
    for cat in grouped:
        content = "\n".join(grouped[cat])
        body += f"\n---\nCategory: {cat}\n{content}\n"
    return header + body

def analyze_pitch(slides):
    prompt = build_analysis_prompt(slides)
    return query_together(prompt, max_tokens=3000)

def parse_analysis_output(raw_text):
    cleaned = raw_text.strip()
    cleaned = re.sub(r"^```json\s*|```$", "", cleaned, flags=re.MULTILINE)
    cleaned = re.sub(r"^```\s*|```$", "", cleaned, flags=re.MULTILINE)
    start = cleaned.find("{")
    end = cleaned.rfind("}")
    if start == -1 or end == -1:
        raise ValueError("Could not find valid JSON object boundaries in LLM output.")
    json_text = cleaned[start:end + 1]
    print(json_text)
    try:
        data = json.loads(json_text)
    except json.JSONDecodeError as e:
        raise ValueError("LLM returned malformed JSON.") from e

    data["processing_date"] = datetime.utcnow().strftime("%d-%m-%Y %H:%M:%S UTC")
    return data


st.set_page_config(page_title="Investment Thesis Generator", layout="centered")

st.title("Automated Investment Thesis Generator")

st.markdown("""
Upload a startup pitch deck in .pptx format. The deck must have 5–20 slides and include at least 3 of the following types:  
Problem, Solution/Product, Market, Business Model, Competition, Team, Financials, Traction, Funding Ask.
""")

uploaded_file = st.file_uploader("Upload your PowerPoint (.pptx only)", type=["pptx"])

if uploaded_file:
    if uploaded_file.size > 50 * 1024 * 1024:
        st.error("File size exceeds 50MB.")
    else:
        prs = Presentation(uploaded_file)
        slide_count = len(prs.slides)
        if not (5 <= slide_count <= 20):
            st.error(f"Slide count {slide_count} is out of bounds (5–20).")
        else:
            if st.button("Generate Investment Report"):
                try:
                    with st.spinner("Extracting text from slides..."):
                        slide_data = extract_slide_text(uploaded_file)
                    st.success("Slide text extracted.")

                    with st.spinner("Classifying slides..."):
                        classified = classify_all_slides(slide_data)
                    st.success("Slides classified.")

                    found = {s["category"] for s in classified if s["category"] in REQUIRED_TYPES}
                    if len(found) < 3:
                        st.error(f"Only {len(found)} required slide types found: {', '.join(found)}")
                        st.stop()

                    with st.spinner("Analyzing pitch deck..."):
                        raw = analyze_pitch(classified)
                        parsed = parse_analysis_output(raw)
                    st.success("Analysis complete.")

                    with st.spinner("Generating PDF report..."):
                        pdf_bytes, pdf_name = generate_pdf_report(parsed)
                    st.success("Report ready.")

                    st.subheader("Recommendation Summary")
                    st.write(f"**{parsed['recommendation']}**")
                    st.write(f"Overall Score: {parsed['overall_score']}")
                    st.write(f"Confidence Score: {parsed['confidence_score']}")
                    st.session_state["pdf_bytes"] = pdf_bytes
                    st.session_state["pdf_name"] = pdf_name

                except Exception as e:
                    st.error(f"Error: {e}")

if "pdf_bytes" in st.session_state and "pdf_name" in st.session_state:
    st.download_button("Download PDF Report",
                       data=st.session_state["pdf_bytes"],
                       file_name=st.session_state["pdf_name"],
                       mime="application/pdf")



